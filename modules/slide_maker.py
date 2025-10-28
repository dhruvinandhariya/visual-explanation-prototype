from pptx import Presentation

def create_ppt(slides_data, outdir):
    prs = Presentation()
    for slide in slides_data:
        layout = prs.slide_layouts[1]
        s = prs.slides.add_slide(layout)
        s.shapes.title.text = slide['title']
        body = s.placeholders[1]
        tf = body.text_frame
        for bullet in slide['bullets']:
            p = tf.add_paragraph()
            p.text = bullet
    ppt_path = f"{outdir}/slides.pptx"
    prs.save(ppt_path)
    return ppt_path
